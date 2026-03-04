#!/usr/bin/env python3
"""Generate an architecture diagram PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
COLOR_INPUT = RGBColor(100, 150, 255)      # Blue
COLOR_SERVICE = RGBColor(76, 175, 80)      # Green
COLOR_DB = RGBColor(255, 152, 0)           # Orange
COLOR_AI = RGBColor(156, 39, 176)          # Purple
COLOR_GATEWAY = RGBColor(244, 67, 54)      # Red
COLOR_FRONTEND = RGBColor(33, 150, 243)    # Bright Blue

def add_box(slide, left, top, width, height, text, color, font_size=12):
    """Add a colored box with text."""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(50, 50, 50)
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = 1  # Middle

    return shape

def add_arrow(slide, x1, y1, x2, y2):
    """Add a simple arrow connector."""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = RGBColor(100, 100, 100)
    connector.line.width = Pt(2)

# ===== SLIDE 1: Title =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 240, 240)

title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
p = title_frame.paragraphs[0]
p.text = "Options Trading AI Platform"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
subtitle_frame = subtitle_box.text_frame
p = subtitle_frame.paragraphs[0]
p.text = "Microservices Architecture\nRaspberry Pi + Cloud"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(100, 100, 100)
p.alignment = PP_ALIGN.CENTER

# ===== SLIDE 2: System Architecture =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

# Title
title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "System Architecture"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)

# Layer 1: Input
add_box(slide2, 3.5, 0.9, 3, 0.6, "Schwab API", COLOR_INPUT, 11)
add_arrow(slide2, 5, 1.5, 5, 2.0)

# Layer 2: Ingestion
add_box(slide2, 2.5, 2.0, 5, 0.6, "Data Ingestion Service", COLOR_SERVICE, 11)
add_arrow(slide2, 5, 2.6, 5, 3.1)

# Layer 3: Database
add_box(slide2, 3.5, 3.1, 3, 0.6, "InfluxDB + Postgres", COLOR_DB, 11)
add_arrow(slide2, 4.5, 3.7, 3.5, 4.2)
add_arrow(slide2, 5.5, 3.7, 6.5, 4.2)

# Layer 4: Processing
add_box(slide2, 0.5, 4.2, 3, 0.6, "Quant Engine", COLOR_SERVICE, 10)
add_box(slide2, 6.5, 4.2, 3, 0.6, "Portfolio Service", COLOR_SERVICE, 10)
add_arrow(slide2, 2, 4.8, 2.5, 5.3)
add_arrow(slide2, 8, 4.8, 7.5, 5.3)

# Layer 5: AI Agents
add_box(slide2, 0.3, 5.3, 2.4, 0.6, "Options AI", COLOR_AI, 10)
add_box(slide2, 3, 5.3, 2.4, 0.6, "Fundamental AI", COLOR_AI, 10)
add_box(slide2, 5.7, 5.3, 2.4, 0.6, "Risk Engine", COLOR_GATEWAY, 10)
add_arrow(slide2, 2, 5.9, 4, 6.4)
add_arrow(slide2, 4, 5.9, 5, 6.4)
add_arrow(slide2, 7, 5.9, 5.5, 6.4)

# Layer 6: Backend
add_box(slide2, 3.5, 6.4, 3, 0.6, "FastAPI Gateway", COLOR_GATEWAY, 11)
add_arrow(slide2, 5, 7.0, 5, 7.5)

# Layer 7: Frontend
add_box(slide2, 3.5, 7.5, 3, 0.6, "React Dashboard", COLOR_FRONTEND, 11)

# ===== SLIDE 3: Microservices =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Microservices"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)

services = [
    ("Ingestion Service", "Schwab API → InfluxDB", COLOR_SERVICE),
    ("Quant Engine", "IV percentile, skew, spreads", COLOR_SERVICE),
    ("Options AI Agent", "LLM: recommend sell strategies", COLOR_AI),
    ("Fundamental Agent", "EDGAR RAG, DCF, cash flow", COLOR_AI),
    ("Portfolio Service", "Real-time P&L, cost basis", COLOR_SERVICE),
    ("Risk Engine", "Margin, exposure validation", COLOR_GATEWAY),
    ("API Gateway", "FastAPI REST + WebSocket", COLOR_GATEWAY),
    ("Worker (Celery)", "Background jobs, EDGAR fetch", COLOR_SERVICE),
]

y_pos = 1.2
for name, desc, color in services:
    add_box(slide3, 0.5, y_pos, 4.2, 0.5, name, color, 10)

    desc_box = slide3.shapes.add_textbox(Inches(5), Inches(y_pos), Inches(4.5), Inches(0.5))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(80, 80, 80)
    tf.vertical_anchor = 1

    y_pos += 0.65

# ===== SLIDE 4: Data Stores =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Data Stores"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)

stores = [
    ("InfluxDB", "Time-series: ticks, IV, Greeks, prices", COLOR_DB),
    ("Postgres", "Recommendations, positions, P&L, audit", COLOR_DB),
    ("pgvector", "EDGAR filings RAG (semantic search)", COLOR_DB),
    ("Redis", "Celery broker, WebSocket pub/sub, cache", COLOR_DB),
]

y_pos = 1.5
for name, desc, color in stores:
    add_box(slide4, 1, y_pos, 3, 0.6, name, color, 11)

    desc_box = slide4.shapes.add_textbox(Inches(4.2), Inches(y_pos), Inches(5.3), Inches(0.6))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(80, 80, 80)
    tf.vertical_anchor = 1

    y_pos += 1.1

# ===== SLIDE 5: Frontend Pages =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Frontend Pages"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)

pages = [
    ("Dashboard", "Live options chain, price charts, WebSocket feed"),
    ("Options Agent", "AI sell ideas: spreads, puts, covered calls"),
    ("Fundamental Agent", "DCF, revenue/earnings charts, EDGAR analysis"),
    ("Portfolio", "Real-time P&L, positions, cost basis tracking"),
    ("Settings", "Ticker management, API keys, agent config"),
]

y_pos = 1.3
for page, desc in pages:
    add_box(slide5, 0.5, y_pos, 2.5, 0.5, page, COLOR_FRONTEND, 10)

    desc_box = slide5.shapes.add_textbox(Inches(3.2), Inches(y_pos), Inches(6.3), Inches(0.5))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(80, 80, 80)
    tf.vertical_anchor = 1

    y_pos += 0.9

# ===== SLIDE 6: Tech Stack =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Technology Stack"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)

# Column 1
col1_title = slide6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(0.4))
tf = col1_title.text_frame
p = tf.paragraphs[0]
p.text = "Backend"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)

col1_items = [
    "Python 3.9+",
    "FastAPI",
    "AsyncIO",
    "Pydantic",
    "SQLAlchemy",
    "Celery",
]
y = 1.7
for item in col1_items:
    box = slide6.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(60, 60, 60)
    y += 0.45

# Column 2
col2_title = slide6.shapes.add_textbox(Inches(3.8), Inches(1.2), Inches(3), Inches(0.4))
tf = col2_title.text_frame
p = tf.paragraphs[0]
p.text = "Frontend"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)

col2_items = [
    "React 18+",
    "Vite",
    "TailwindCSS",
    "WebSocket",
    "TypeScript",
]
y = 1.7
for item in col2_items:
    box = slide6.shapes.add_textbox(Inches(3.8), Inches(y), Inches(3), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(60, 60, 60)
    y += 0.45

# Column 3
col3_title = slide6.shapes.add_textbox(Inches(7.1), Inches(1.2), Inches(2.4), Inches(0.4))
tf = col3_title.text_frame
p = tf.paragraphs[0]
p.text = "AI / Data"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(33, 150, 243)

col3_items = [
    "Claude API",
    "InfluxDB",
    "PostgreSQL",
    "pgvector",
    "Redis",
]
y = 1.7
for item in col3_items:
    box = slide6.shapes.add_textbox(Inches(7.1), Inches(y), Inches(2.4), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(60, 60, 60)
    y += 0.45

# Save
prs.save('/home/umahar/stocks/Architecture.pptx')
print("✓ Architecture.pptx created successfully!")
